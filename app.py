#!/usr/bin/env python3
"""
Claude Web Chat v2.0 - Backend Server
Features: Multi-file upload, Document generation (Word/Excel/PPT), Skills, SSH
"""

from flask import Flask, request, jsonify, send_from_directory, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import os
import io
import re
import json
import base64
import paramiko
import requests
from pathlib import Path
from datetime import datetime

from dotenv import load_dotenv
load_dotenv(dotenv_path='env')

app = Flask(__name__, static_folder='.', static_url_path='')
CORS(app)

API_KEY      = os.getenv("ANTHROPIC_API_KEY")
SSH_HOST     = os.getenv("SSH_HOST", "")
SSH_PORT     = int(os.getenv("SSH_PORT", "22"))
SSH_USER     = os.getenv("SSH_USER", "")
SSH_PASSWORD = os.getenv("SSH_PASSWORD", "")
SSH_KEY_PATH = os.getenv("SSH_KEY_PATH", "")

HISTORY_DIR = Path("./chat_history")
HISTORY_DIR.mkdir(exist_ok=True)

ANTHROPIC_API_URL = "https://api.anthropic.com/v1/messages"
MAX_FILE_TEXT     = 15000  # max chars extracted from a file

IMAGE_EXTS = {'jpg', 'jpeg', 'png', 'gif', 'webp', 'bmp'}
MEDIA_TYPES = {
    'jpg': 'image/jpeg', 'jpeg': 'image/jpeg', 'png': 'image/png',
    'gif': 'image/gif',  'webp': 'image/webp',  'bmp': 'image/bmp',
}

# ============================================================================
# CLAUDE CHAT HANDLER
# ============================================================================

class ClaudeChat:
    def __init__(self):
        self.ssh_client = None

    def get_system_prompt(self):
        return """Bạn là Claude - AI Assistant toàn năng với nhiều kỹ năng chuyên biệt.

## SKILLS CÓ SẴN:

### 📄 TẠO TÀI LIỆU
- **Word/Docx**: Báo cáo, hợp đồng, đề xuất. Dùng `# Tiêu đề`, `## Mục`, `- bullet`, `**bold**`
- **Excel**: Bảng dữ liệu, thống kê. Dùng format bảng markdown `| cột | cột |` để tạo sheet đẹp
- **PowerPoint**: Dùng `## Tên Slide` cho từng slide, `- nội dung` cho bullet points

### 💻 LẬP TRÌNH
- **Web Frontend**: HTML5, CSS3, JavaScript, React, Vue, Tailwind CSS — code hoàn chỉnh, chạy được ngay
- **Web Backend**: Python/Flask/Django, Node.js/Express, REST API, Database schema
- **Mobile App**: React Native, Flutter, Kotlin, Swift
- **DevOps/API**: Docker, CI/CD, REST API design, OpenAPI spec

### 📊 PHÂN TÍCH & NGHIÊN CỨU
- **Data Analysis**: Phân tích với pandas/numpy, tạo code visualization (matplotlib/plotly)
- **Scientific Research**: Tổng hợp tài liệu, phương pháp luận, review tài liệu khoa học
- **Statistics**: Mô tả thống kê, kiểm định giả thuyết, hồi quy, phân tích tương quan

### 🛠️ KHÁC
- **SQL**: Query phức tạp, tối ưu, thiết kế schema
- **Translation**: Dịch đa ngôn ngữ, giữ nguyên văn phong
- **Content Writing**: Blog, báo cáo, email, marketing
- **SSH Execution**: Thực thi lệnh trên server từ xa

## QUY TẮC:
1. Khi user upload file → đọc và phân tích nội dung file, trả lời dựa trên đó
2. Khi tạo Word/Excel/PPT → dùng markdown chuẩn, user nhấn nút tải về để convert
3. Khi code → cung cấp code đầy đủ, có thể chạy ngay, có giải thích
4. Trả lời tiếng Việt nếu user dùng tiếng Việt
5. Chia output thành sections rõ ràng với heading markdown"""

    def chat(self, messages, model="claude-sonnet-4-20250514"):
        try:
            if not API_KEY:
                return "❌ Error: ANTHROPIC_API_KEY not configured"

            headers = {
                "x-api-key": API_KEY,
                "anthropic-version": "2023-06-01",
                "content-type": "application/json",
            }
            payload = {
                "model": model,
                "max_tokens": 4096,
                "system": self.get_system_prompt(),
                "messages": messages,
            }
            response = requests.post(ANTHROPIC_API_URL, json=payload, headers=headers, timeout=60)
            if response.status_code != 200:
                return "❌ API Error: {}".format(response.text[:200])
            return response.json()['content'][0]['text']

        except Exception as e:
            return "❌ Error: {}".format(str(e))

    # ---- SSH ----------------------------------------------------------------

    def _is_ssh_connected(self):
        try:
            if self.ssh_client is None:
                return False
            transport = self.ssh_client.get_transport()
            return transport is not None and transport.is_active()
        except Exception:
            return False

    def ssh_execute(self, command):
        try:
            if not SSH_HOST or not SSH_USER:
                return "❌ SSH not configured"

            if not self._is_ssh_connected():
                self.ssh_client = paramiko.SSHClient()
                self.ssh_client.set_missing_host_key_policy(paramiko.WarningPolicy())
                if SSH_KEY_PATH and os.path.exists(SSH_KEY_PATH):
                    self.ssh_client.connect(SSH_HOST, port=SSH_PORT, username=SSH_USER,
                                            key_filename=SSH_KEY_PATH, timeout=10)
                else:
                    self.ssh_client.connect(SSH_HOST, port=SSH_PORT, username=SSH_USER,
                                            password=SSH_PASSWORD, timeout=10)

            stdin, stdout, stderr = self.ssh_client.exec_command(command, timeout=30)
            output = stdout.read().decode('utf-8', errors='ignore')
            error  = stderr.read().decode('utf-8', errors='ignore')
            if output and error:
                return output + "\n" + error
            return output if output else error

        except Exception as e:
            self.ssh_client = None
            return "❌ SSH Error: {}".format(str(e))

    def test_ssh(self):
        result = self.ssh_execute("whoami")
        if "❌" in result:
            return {"success": False, "message": "SSH failed: {}".format(result)}
        return {"success": True, "message": "SSH OK - User: {}".format(result.strip())}


claude_chat = ClaudeChat()

# ============================================================================
# FILE PROCESSING
# ============================================================================

def process_uploaded_file(raw_bytes, filename):
    """Return a dict describing the file for the AI."""
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    # ---- Images → base64 for Claude Vision ----
    if ext in IMAGE_EXTS:
        return {
            'type': 'image',
            'name': filename,
            'data': base64.b64encode(raw_bytes).decode('utf-8'),
            'mediaType': MEDIA_TYPES.get(ext, 'image/jpeg'),
        }

    # ---- PDF ----
    if ext == 'pdf':
        try:
            import pypdf
            reader = pypdf.PdfReader(io.BytesIO(raw_bytes))
            text = '\n'.join(p.extract_text() or '' for p in reader.pages)
            return {'type': 'text', 'name': filename, 'content': text[:MAX_FILE_TEXT]}
        except Exception as e:
            return {'type': 'text', 'name': filename, 'content': '[PDF lỗi: {}]'.format(e)}

    # ---- Word ----
    if ext == 'docx':
        try:
            from docx import Document
            doc = Document(io.BytesIO(raw_bytes))
            text = '\n'.join(p.text for p in doc.paragraphs if p.text)
            return {'type': 'text', 'name': filename, 'content': text[:MAX_FILE_TEXT]}
        except Exception as e:
            return {'type': 'text', 'name': filename, 'content': '[DOCX lỗi: {}]'.format(e)}

    # ---- Excel ----
    if ext in ('xlsx', 'xls'):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
            lines = []
            for sheet in wb.worksheets:
                lines.append('=== Sheet: {} ==='.format(sheet.title))
                for row in sheet.iter_rows(values_only=True):
                    lines.append('\t'.join('' if c is None else str(c) for c in row))
            return {'type': 'text', 'name': filename, 'content': '\n'.join(lines)[:MAX_FILE_TEXT]}
        except Exception as e:
            return {'type': 'text', 'name': filename, 'content': '[Excel lỗi: {}]'.format(e)}

    # ---- CSV ----
    if ext == 'csv':
        text = raw_bytes.decode('utf-8', errors='ignore')
        return {'type': 'text', 'name': filename, 'content': text[:MAX_FILE_TEXT]}

    # ---- Video / Audio (unsupported by Claude) ----
    if ext in ('mp4', 'avi', 'mov', 'mkv', 'webm', 'wmv', 'flv'):
        return {'type': 'text', 'name': filename,
                'content': '[Video: {} — AI không thể xử lý video trực tiếp]'.format(filename)}
    if ext in ('mp3', 'wav', 'ogg', 'flac', 'aac', 'm4a'):
        return {'type': 'text', 'name': filename,
                'content': '[Audio: {} — AI không thể xử lý audio trực tiếp]'.format(filename)}

    # ---- Fallback: text / source code ----
    try:
        text = raw_bytes.decode('utf-8', errors='ignore')
        return {'type': 'text', 'name': filename, 'content': text[:MAX_FILE_TEXT]}
    except Exception:
        return {'type': 'text', 'name': filename, 'content': '[Binary file: {}]'.format(filename)}


# ============================================================================
# DOCUMENT GENERATION HELPERS
# ============================================================================

def _build_word(content, title):
    from docx import Document as DocxDoc
    doc = DocxDoc()
    doc.add_heading(title, 0)
    for line in content.split('\n'):
        if line.startswith('### '):
            doc.add_heading(line[4:], level=3)
        elif line.startswith('## '):
            doc.add_heading(line[3:], level=2)
        elif line.startswith('# '):
            doc.add_heading(line[2:], level=1)
        elif line.startswith(('- ', '* ')):
            doc.add_paragraph(line[2:], style='List Bullet')
        elif line.strip():
            parts = re.split(r'\*\*(.+?)\*\*', line)
            p = doc.add_paragraph()
            for i, part in enumerate(parts):
                run = p.add_run(part)
                if i % 2 == 1:
                    run.bold = True
        else:
            doc.add_paragraph()
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _build_excel(content, title):
    import openpyxl
    from openpyxl.styles import Font
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = title[:31]
    row_num = 1
    for line in content.split('\n'):
        if '|' in line and line.strip().startswith('|'):
            cells = [c.strip() for c in line.split('|')[1:-1]]
            if all(set(c) <= set('- :') for c in cells if c):
                continue  # markdown separator row
            for col, val in enumerate(cells, 1):
                cell = ws.cell(row=row_num, column=col, value=val)
                if row_num == 1:
                    cell.font = Font(bold=True)
            row_num += 1
        elif line.strip():
            ws.cell(row=row_num, column=1, value=line.strip())
            row_num += 1
    for col in ws.columns:
        max_len = max((len(str(c.value)) for c in col if c.value), default=10)
        ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 60)
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _build_ppt(content, title):
    from pptx import Presentation
    prs = Presentation()

    # Title slide
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = title
    try:
        s.placeholders[1].text = 'Tạo bởi Claude AI'
    except Exception:
        pass

    cur_title, cur_bullets = None, []

    def flush():
        if cur_title is None:
            return
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = cur_title
        if cur_bullets:
            try:
                tf = sl.placeholders[1].text_frame
                tf.clear()
                for i, b in enumerate(cur_bullets):
                    if i == 0:
                        tf.paragraphs[0].text = b
                    else:
                        tf.add_paragraph().text = b
            except Exception:
                pass

    for line in content.split('\n'):
        if line.startswith('## ') or line.startswith('# '):
            flush()
            cur_title   = line.lstrip('# ').strip()
            cur_bullets = []
        elif line.startswith('- ') or line.startswith('* '):
            cur_bullets.append(line[2:])
        elif line.startswith('### '):
            cur_bullets.append(line[4:])
        elif line.strip() and cur_title is not None:
            cur_bullets.append(line.strip())

    flush()
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ============================================================================
# ROUTES
# ============================================================================

@app.route('/')
def index():
    return send_from_directory('.', 'index.html')


@app.route('/api/health', methods=['GET'])
def health():
    return jsonify({
        "status": "ok", "version": "2.0.0",
        "api_configured": bool(API_KEY),
        "timestamp": datetime.now().isoformat()
    })


@app.route('/api/chat', methods=['POST'])
def chat_endpoint():
    try:
        data     = request.json
        messages = data.get('messages', [])
        model    = data.get('model', 'claude-sonnet-4-20250514')
        if not messages:
            return jsonify({"error": "No messages provided"}), 400
        response = claude_chat.chat(messages, model)
        return jsonify({"success": True, "response": response,
                        "timestamp": datetime.now().isoformat()})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


# ---- File upload ------------------------------------------------------------

@app.route('/api/upload', methods=['POST'])
def upload_files():
    try:
        if 'files' not in request.files:
            return jsonify({'error': 'No files provided'}), 400
        results = []
        for f in request.files.getlist('files'):
            if not f.filename:
                continue
            filename = secure_filename(f.filename)
            result   = process_uploaded_file(f.read(), filename)
            results.append(result)
        return jsonify({'success': True, 'files': results})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


# ---- Document generation ---------------------------------------------------

@app.route('/api/generate/word', methods=['POST'])
def generate_word():
    try:
        data    = request.json
        content = data.get('content', '')
        title   = data.get('title', 'Document')
        buf     = _build_word(content, title)
        return send_file(buf,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True, download_name=title + '.docx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate/excel', methods=['POST'])
def generate_excel():
    try:
        data    = request.json
        content = data.get('content', '')
        title   = data.get('title', 'Spreadsheet')
        buf     = _build_excel(content, title)
        return send_file(buf,
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            as_attachment=True, download_name=title + '.xlsx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/generate/ppt', methods=['POST'])
def generate_ppt():
    try:
        data    = request.json
        content = data.get('content', '')
        title   = data.get('title', 'Presentation')
        buf     = _build_ppt(content, title)
        return send_file(buf,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True, download_name=title + '.pptx')
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ---- SSH -------------------------------------------------------------------

@app.route('/api/ssh/test', methods=['GET'])
def ssh_test():
    return jsonify(claude_chat.test_ssh())


@app.route('/api/ssh/execute', methods=['POST'])
def ssh_execute():
    try:
        data    = request.json
        command = data.get('command')
        if not command:
            return jsonify({"error": "No command provided"}), 400
        result = claude_chat.ssh_execute(command)
        return jsonify({"success": True, "command": command, "result": result,
                        "timestamp": datetime.now().isoformat()})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


# ---- History ---------------------------------------------------------------

@app.route('/api/history/save', methods=['POST'])
def save_history():
    try:
        data       = request.json
        session_id = data.get('session_id', 'default')
        messages   = data.get('messages', [])
        # Strip large image data before saving to keep history files small
        slim = []
        for m in messages:
            if isinstance(m.get('content'), list):
                slim_content = []
                for block in m['content']:
                    if block.get('type') == 'image':
                        slim_content.append({'type': 'text', 'text': '[Image: skipped in history]'})
                    else:
                        slim_content.append(block)
                slim.append({'role': m['role'], 'content': slim_content})
            else:
                slim.append(m)
        history_file = HISTORY_DIR / "{}.json".format(session_id)
        with open(str(history_file), 'w', encoding='utf-8') as f:
            json.dump({"session_id": session_id, "messages": slim,
                       "saved_at": datetime.now().isoformat()}, f,
                      ensure_ascii=False, indent=2)
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/history/load/<session_id>', methods=['GET'])
def load_history(session_id):
    try:
        history_file = HISTORY_DIR / "{}.json".format(session_id)
        if not history_file.exists():
            return jsonify({"success": False, "error": "Session not found"}), 404
        with open(str(history_file), 'r', encoding='utf-8') as f:
            data = json.load(f)
        return jsonify({"success": True, "data": data})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/history/list', methods=['GET'])
def list_histories():
    try:
        histories = []
        for file in HISTORY_DIR.glob("*.json"):
            with open(str(file), 'r', encoding='utf-8') as f:
                d = json.load(f)
            histories.append({
                "session_id":    d.get("session_id"),
                "saved_at":      d.get("saved_at"),
                "message_count": len(d.get("messages", []))
            })
        return jsonify({"success": True,
                        "histories": sorted(histories, key=lambda x: x['saved_at'], reverse=True)})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


# ============================================================================
# ERROR HANDLERS
# ============================================================================

@app.errorhandler(404)
def not_found(error):
    return jsonify({"error": "Not found"}), 404

@app.errorhandler(500)
def server_error(error):
    return jsonify({"error": "Internal server error"}), 500


# ============================================================================
# MAIN
# ============================================================================

if __name__ == '__main__':
    print("\n" + "="*60)
    print("  Claude Web Chat v2.0")
    print("="*60)
    print("  http://localhost:5000")
    print("="*60 + "\n")
    app.run(debug=False, host='0.0.0.0', port=5000, threaded=True)
